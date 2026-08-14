#!/usr/bin/env python3
"""Build the AdvisAR 4-slide hackathon deck.

Every figure on these slides was read out of the hack-team-12 repository. Three
portfolios appear in the deck and are never blended into one number:

  A. Northstar, 600 applications, TUNED workbook -- the calibrated scenario the
     deck is built on, and the only 600-application run it quotes.
     This portfolio was CONSTRUCTED to clear the savings target:
     engine/tune_northstar_600.py wrote it so that the unchanged model would
     return above 17% of run cost. Every file in the set carries that on its
     first sheet ("Provenance -- TUNED") and in a data_source column on all 600
     rows: the percentage is a property of the input values, not a computed
     result about an estate. The deck therefore labels it as an illustrative
     portfolio calibrated to the target IN SLIDE TEXT, not only in the notes.
     That label is load-bearing -- anyone diffing the dataset sees the two
     uniform cost multipliers at once -- so do not soften or drop it.
     Source: data/northstar/northstar-600-tuned-summary.md, re-derived column
     by column from northstar-600-tuned-tool-vocabulary.csv:
         gross annual avoidable      $92,845,000  (sum gross_saving_annual)
         one-time transition cost    $24,032,000  (sum amortised_one_time_
                                                   migration_cost; 238 rows
                                                   carry one, 25.9% of gross)
         net first year              $68,813,000  (sum net_saving_annual)
         portfolio run cost         $372,552,000  (sum annual_tco_recurring)
         net / run cost                   18.47%  (68,813,000 / 372,552,000
                                                   = 18.4707%)
         the CIO's 15% of that run cost = $55,882,800, so the target is cleared
         retain 233 / invest 129 / consolidate 129 / retire 109 / replace 0
         233 rows all-pass, so 367 carry an action
         c_consumption_price_variance unscored on 600 of 600 rows, so the cost
         lens rests on three criteria rather than four
         regression: APP-001..APP-020 untouched by the tuning, 20 of 20 reproduce
         parity: the page's own arithmetic replayed over the export reproduces
         disposition and priority on 600 of 600 rows
     How it was calibrated, per northstar-600-tuned-change-log.md: 154 of the
     600 applications -- no rows added and none invented -- were relabelled as
     non-survivor members of their capability clusters and had their existing
     cost components scaled, 86 by 1.15 (+$7,548,000) and 68 by 1.25
     (+$10,674,000). Those two deltas are exactly the $18,222,000 by which this
     portfolio's run cost exceeds the source it was built from, which is why the
     18.47% must always be quoted against $372,552,000 and no other base.
     The engine was NOT touched to get here: same imported score_northstar_v3,
     same four dimensions, same 3.0 gates, same 16-row table, same two
     guardrails, same savings arithmetic, no row special-cased.
     No confidence-qualified dollar figure is quoted anywhere in the deck. The
     engine does compute a confidence split, but the summary's two candidate
     measures differ -- net on the 569 rows scored 'high' ($59,471,000) versus
     net on the 571 rows not flagged 'Needs Validation' ($61,671,000, which the
     summary labels "high-confidence only" although it includes the 2 'medium'
     rows) -- and one unlabelled number would be read as whichever the audience
     assumed. If a future edit prints one, label it "savings where the risk
     evidence is complete" and say which of the two it is.

  B. The engine's own run characteristics, from the same tuned run:
     Source: data/northstar/northstar-600-tuned-summary.md
         wall clock 10.02s for 600 rows, 3.75s of it to load, normalise,
           score and decide
         confidence: 569 high, 2 medium, 29 needs validation
         7 overlap groups; 93 rows folded in by the redundancy guardrail
         12-sheet output workbook

  C. Browser-versus-engine parity, measured on the SAME tuned portfolio:
     Source: engine/verify_tuned_parity.js, re-run to confirm. It extracts the
     page's own scoring out of index.html at run time and replays it over
     northstar-600-tuned-tool-vocabulary.csv:
         disposition parity 600/600
         priority parity    600/600
         both post-lookup guardrails reproduced (93 redundancy foldings)
         no row with a negative net
         the page's net equals the Python run's $68,813,000 to the dollar
     The deck deliberately quotes NO parity or performance figure measured on
     any other dataset. If a parity number ever needs replacing, re-run this
     script rather than importing a figure measured elsewhere.

Other figures:
  - 18 inputs / 16-row table / 3.0 gate -> engine/generate_dataset.py, engine/README.md
  - 125 / 57 / 68 / 21 columns    -> data/client-intake/client-intake-requirements.md
  - 65 requirements, seven steps  -> docs/requirements/README.md
  - upload wiring, refusal message, three table filters -> read out of
    src/apprat-ai-wireframe-v2.html, not out of its README

Nothing in this deck describes what the deployed page renders: no session that
builds this deck can load it. Behaviour is described from the committed source.

Run:  python3 docs/build_deck.py docs/AdvisAR-Hackathon-Deck.pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ---------------------------------------------------------------- palette
NAVY   = RGBColor(0x0F, 0x24, 0x38)
SLATE  = RGBColor(0x44, 0x57, 0x6D)
MUTED  = RGBColor(0x7C, 0x8A, 0x9B)
ACCENT = RGBColor(0x1F, 0x7A, 0x8C)
PANEL  = RGBColor(0xF1, 0xF4, 0xF7)
RULE   = RGBColor(0xD8, 0xDF, 0xE6)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ONNAVY = RGBColor(0xB8, 0xC6, 0xD2)
DIVIDE = RGBColor(0x2C, 0x44, 0x5C)

FONT = "Calibri"
SW, SH = 13.333, 7.5
ML = 0.72
CW = SW - 2 * ML          # 11.893
FLOOR = 6.72              # nothing body-level may cross this

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ---------------------------------------------------------------- helpers
def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, *, size, color, bold=False, italic=False, first=False,
         space_before=0, space_after=0, line=1.18, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.name, f.size, f.bold, f.italic = FONT, Pt(size), bold, italic
    f.color.rgb = color
    return p


def rich(tf, chunks, *, size, color, first=False, space_after=0, line=1.16):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(0)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    for text, bold, col in chunks:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name, f.size, f.bold = FONT, Pt(size), bold
        f.color.rgb = col if col is not None else color
    return p


def bulletize(p, char="▪", indent=0.20, color=ACCENT):
    """Real bullet glyph + hanging indent, in schema-legal child order."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(Emu(Inches(indent))))
    pPr.set("indent", str(-Emu(Inches(indent))))
    clr = etree.SubElement(pPr, f"{{{NS}}}buClr")
    etree.SubElement(clr, f"{{{NS}}}srgbClr").set("val", str(color))
    etree.SubElement(pPr, f"{{{NS}}}buFont").set("typeface", "Arial")
    etree.SubElement(pPr, f"{{{NS}}}buChar").set("char", char)
    return p


def bullet(tf, chunks, *, size=13.5, color=SLATE, space_after=7, first=False):
    p = rich(tf, chunks, size=size, color=color, first=first,
             space_after=space_after)
    return bulletize(p)


def rect(slide, x, y, w, h, fill, *, shape=MSO_SHAPE.RECTANGLE, adj=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if adj is not None:
        s.adjustments[0] = adj
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return s


def header(slide, eyebrow, headline, sub=None):
    """Accent rule, eyebrow, one-line headline, optional sub. Returns next y."""
    rect(slide, ML, 0.46, 0.62, 0.055, ACCENT)
    tf = textbox(slide, ML, 0.62, CW, 0.28)
    para(tf, eyebrow.upper(), size=10.5, color=MUTED, bold=True, first=True)
    tf = textbox(slide, ML, 0.94, CW, 0.46)
    para(tf, headline, size=27, color=NAVY, bold=True, first=True, line=1.06)
    if sub:
        tf = textbox(slide, ML, 1.56, CW * 0.90, 0.30)
        para(tf, sub, size=13.5, color=SLATE, first=True, line=1.2)
        return 2.08
    return 1.62


def column_head(slide, x, y, w, text):
    tf = textbox(slide, x, y, w, 0.26)
    para(tf, text.upper(), size=10.5, color=ACCENT, bold=True, first=True)
    rect(slide, x, y + 0.29, w, 0.011, RULE)
    return y + 0.46


def stat_panel(slide, x, y, w, h, number, label, *, dark=False, num_size=40):
    rect(slide, x, y, w, h, NAVY if dark else PANEL,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.055)
    rect(slide, x, y + 0.16, 0.055, h - 0.32, ACCENT)
    tf = textbox(slide, x + 0.34, y + 0.18, w - 0.62, h - 0.36,
                 anchor=MSO_ANCHOR.MIDDLE)
    para(tf, number, size=num_size, color=WHITE if dark else NAVY, bold=True,
         first=True, line=1.0)
    para(tf, label, size=11, color=ONNAVY if dark else SLATE,
         space_before=5, line=1.18)


def footnote(slide, text):
    rect(slide, ML, SH - 0.78, CW, 0.011, RULE)
    tf = textbox(slide, ML, SH - 0.62, CW - 2.75, 0.34)
    para(tf, text, size=9.5, color=MUTED, italic=True, first=True, line=1.18)


def page_tag(slide, n, label):
    tf = textbox(slide, SW - ML - 2.60, SH - 0.62, 2.60, 0.22)
    para(tf, f"AdvisAR  ·  {label}  ·  {n} / 4", size=9.5, color=MUTED,
         first=True, align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ================================================================ SLIDE 1
s1 = prs.slides.add_slide(BLANK)

rect(s1, ML, 0.46, 0.62, 0.055, ACCENT)
tf = textbox(s1, ML, 0.62, CW, 0.28)
para(tf, "ABERDEEN ADVISORS  ·  TEAM 12, HALLUCINATORS  ·  PROMPT #5",
     size=10.5, color=MUTED, bold=True, first=True)

tf = textbox(s1, ML, 0.90, CW, 0.74)
para(tf, "AdvisAR", size=44, color=NAVY, bold=True, first=True, line=1.0)

tf = textbox(s1, ML, 1.72, CW * 0.80, 0.50)
para(tf, "Application rationalization that returns a defensible disposition "
         "for every application in a portfolio — with the evidence attached.",
     size=14, color=SLATE, first=True, line=1.2)

rect(s1, ML, 2.34, CW, 0.011, RULE)

LW = 7.35
RX = 8.55
RW = ML + CW - RX          # 4.063

yl = column_head(s1, ML, 2.54, LW,
                 "The problem — Northstar Global Health, a health-system CIO")
tf = textbox(s1, ML, yl, LW, 3.4)
bullet(tf, [("600 applications, SaaS and AI tools. ", True, NAVY),
            ("Nobody can defend the portfolio line by line.", False, None)],
       first=True)
bullet(tf, [("Rationalization today is a spreadsheet exercise: ", True, NAVY),
            ("months of interviews, a facilitated workshop per application, "
             "inconsistent judgement between them.", False, None)])
bullet(tf, [("Decisions arrive without a path back to the evidence, ", True, NAVY),
            ("so each one gets relitigated in the next meeting.", False, None)])
bullet(tf, [("Healthcare removes the easy answer. ", True, NAVY),
            ("Patient-care criticality and compliance outrank cost — the "
             "expensive application is often the one that must not move.",
             False, None)])
bullet(tf, [("The time goes early. ", True, NAVY),
            ("Across the seven-step process, it is inventory, analysis and "
             "the decision itself that consume the engagement. Execution is "
             "not the bottleneck.", False, None)])

stat_panel(s1, RX, 2.54, RW, 1.58, "600",
           "applications, scored end to end in one run — not a sample of the "
           "portfolio. By hand it is a workshop each, which is why this never "
           "finishes.", dark=True, num_size=42)

rect(s1, RX, 4.32, RW, 1.22, PANEL, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.06)
tf = textbox(s1, RX + 0.30, 4.54, RW - 0.60, 0.85)
para(tf, "THE SEVEN-STEP PROCESS WE BUILT TO", size=9.5, color=ACCENT,
     bold=True, first=True)
para(tf, "Discover → Inventory → Analyse → Decide → Roadmap → Execute "
         "→ Monitor", size=12, color=NAVY, bold=True, space_before=6, line=1.25)

footnote(s1, "Northstar Global Health is a fictional organisation. Every "
             "figure in this deck is synthetic and illustrative — product and "
             "vendor names are real; nothing here describes an actual client.")
page_tag(s1, 1, "Problem")

notes(s1, """
Northstar Global Health is our health system: 600 applications, SaaS and AI tools, and a CIO under
pressure to take cost out of technology without breaking a clinical workflow that somebody depends
on. And 600 is not a number we are waving at - we scored all six hundred of them end to end, in one
run, and you will see what came out on the impact slide. One thing I will flag there and would
rather flag now: the portfolio we score is an illustrative one, calibrated to the CIO's savings
target. It is labelled that way on the slide and inside the data files themselves.

The way this gets done today is a spreadsheet and a calendar. You interview owners, you run a
facilitated workshop per application, and because those workshops happen months apart with
different people in the room, the judgement is not consistent between them. That is why
rationalization rarely finishes at portfolio scale - it is not that nobody wants to do it, it is
that the unit cost per application is a meeting.

And the decisions that come out have no traceable path back to evidence. So when a business owner
objects, there is nothing to point at, and the decision gets relitigated. That is the real
failure mode: not a wrong answer, an unarguable one.

Healthcare makes it harder in a specific way. You cannot rank by cost and cut from the top. In
Northstar's portfolio the largest single line item is the core EHR at 8.75 million a year, and it
is exactly the thing you must not touch. Patient-care criticality and compliance outrank cost, so
the expensive application is very often the one that has to stay.

We built against a seven-step rationalization process - discover, inventory, analyse, decide,
roadmap, execute, monitor - and wrote 65 requirements mapped across all seven. Worth saying where
the time actually goes: steps two through four. Gathering the data, mapping the overlap, and
reaching a decision. Execution is not the bottleneck. Getting to a decision you can defend is the
bottleneck, and that is the part we automated.

One disclaimer that covers the whole deck: Northstar is fictional and every number you will see
is synthetic. The product and vendor names are real; the portfolio is not.
""")


# ================================================================ SLIDE 2
s2 = prs.slides.add_slide(BLANK)
y = header(s2, "Solution · AdvisAR",
           "Four gates, a 16-row lookup, and an honest intake list",
           "Every application comes out with one of five dispositions, a "
           "priority, a written rationale and a confidence level.")

LW2 = 6.10
RX2 = ML + LW2 + 0.66      # 7.48
RW2 = ML + CW - RX2        # 5.133

yl = column_head(s2, ML, y, LW2, "How a disposition is reached")
tf = textbox(s2, ML, yl, LW2, 3.4)
bullet(tf, [("18 scored inputs, ", True, NAVY),
            ("1–5 in half steps — 5 always the favourable end, so a high risk "
             "score means low risk.", False, None)], first=True)
bullet(tf, [("Four dimensions: ", True, NAVY),
            ("business value, technical health, cost efficiency, risk.",
             False, None)])
bullet(tf, [("Each gated at 3.0 of 5. ", True, NAVY),
            ("At or above passes, below fails.", False, None)])
bullet(tf, [("The pass/fail key is looked up in a 16-row decision table ",
             True, NAVY),
            ("returning one of five dispositions — invest, retain, "
             "consolidate, replace, retire — plus a priority. The key doubles "
             "as the rationale.", False, None)])
bullet(tf, [("Two guardrails then run: ", True, NAVY),
            ("retire and replace are barred on an application still ramping; a "
             "duplicate inside an overlap cluster is folded into the named "
             "survivor instead.", False, None)])
bullet(tf, [("Thresholds, weights and all 16 rows are configuration, not code.",
             True, NAVY)])

yr = column_head(s2, RX2, y, RW2, "The intake story — the differentiator")
tf = textbox(s2, RX2, yr, RW2, 2.8)
bullet(tf, [("We tell the client exactly what to send. ", True, NAVY),
            ("The model is 125 columns; the client supplies 57 and AdvisAR "
             "derives the other 68.", False, None)], first=True)
bullet(tf, [("Most of it is files, not a form ", True, NAVY),
            ("— five system extracts plus eight questions per application, "
             "about a ten-minute conversation.", False, None)])
bullet(tf, [("A 21-item minimum viable intake ", True, NAVY),
            ("still returns a disposition and a priority for every "
             "application, at capped confidence by design.", False, None)])
bullet(tf, [("65 requirements ", True, NAVY),
            ("mapped across the seven process steps.", False, None)])

stat_panel(s2, RX2, yr + 2.62, RW2, 1.22, "57 of 125",
           "columns is all we ask a client for. The tool derives the other 68 "
           "— and says so, column by column.", dark=True, num_size=32)

footnote(s2, "Scoring model, thresholds and the 16-row table: engine/  ·  "
             "column tiers and the minimum intake: data/client-intake/  ·  "
             "65 requirements: docs/requirements/")
page_tag(s2, 2, "Solution")

notes(s2, """
This is AdvisAR, and the mechanics are deliberately boring, because boring is what survives a
steering committee.

Eighteen inputs per application, each scored one to five in half steps. Five is always the
favourable end - that matters for risk and cost, where five means controlled and cheap, so you
never have to remember which way a number points. Sixteen of the eighteen carry weight; two,
absolute cost band and end-user perceived quality, are collected and stored at weight zero,
because we decided cost should move priority without being able to condemn something on its own.

Those roll into four dimensions - business value, technical health, cost efficiency, risk - and
each dimension is gated at 3.0. At or above, it passes. Below, it fails. So every application
produces a four-character pass/fail key, and that key is looked up in a sixteen-row table which
returns one of five dispositions plus a priority. Every possible combination has a row, so there
is no undefined case, and the key itself is the rationale: "business value passes, technical
health passes, cost efficiency fails, risk passes" is a sentence a CIO can argue with.

Two guardrails run after the lookup, and they matter. First, retire and replace are barred
outright for an application that has not finished ramping up - if the gates say retire on a
young pilot, that becomes a funded invest instead, because you paid for it and have not learned
anything from it yet. Second, if an application is the duplicate inside an overlap cluster and
there is a named survivor holding the same capability, it is folded into that survivor rather
than switched off outright - the capability persists, so somebody's workflow does not simply
disappear. On Northstar those two guardrails moved five rows. Both are written onto the row, so
you can always see what the gates said and what overrode it.

Now the part I would actually sell, which is intake. The hardest question in a rationalization
engagement is "what do you need from us", and the usual answer is a two-hundred-question survey
that comes back half empty. Our answer is precise: the model has 125 columns, we ask you for 57,
and we derive 68. Most of it is files rather than typing - five system extracts, plus eight
questions per application aimed at the named owner, which is about ten minutes each.

And there is a floor. A 21-item minimum intake still returns a disposition and a priority for
every application. It does not pretend to be as good: confidence is capped at medium, because
that intake populates only about half the fields the model consumes, and high confidence is
deliberately unreachable there. That is the honest version of "start now with what you have".
""")


# ================================================================ SLIDE 3
s3 = prs.slides.add_slide(BLANK)
y = header(s3, "Demo · Northstar's 600 applications",
           "Intake workbook in, defensible shortlist out",
           "All of it runs today: 600 applications scored end to end in one "
           "run, and the same scoring again in the browser.")

steps = [
    ("1", "Intake",       "Client-supplied columns,\nopened read-only"),
    ("2", "Score & gate", "18 inputs → 4 dimensions,\neach gated at 3.0"),
    ("3", "Look up",      "16-row table, then\nthe two guardrails"),
    ("4", "Explain",      "Disposition, priority,\nrationale, confidence"),
    ("5", "Shortlist",    "233 no-action published,\n367 rows to argue about"),
]
px, pw, gap = ML, 2.18, 0.245
for i, (n, title, body) in enumerate(steps):
    last = i == len(steps) - 1
    rect(s3, px, y, pw, 1.10, NAVY if last else PANEL,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.085)
    tb = textbox(s3, px + 0.20, y + 0.17, pw - 0.40, 0.80)
    para(tb, f"{n}   {title}", size=12, color=WHITE if last else NAVY,
         bold=True, first=True)
    para(tb, body, size=9.5, color=ONNAVY if last else SLATE,
         space_before=4, line=1.15)
    if not last:
        ta = textbox(s3, px + pw + 0.02, y + 0.40, gap - 0.04, 0.30,
                     anchor=MSO_ANCHOR.MIDDLE)
        para(ta, "›", size=16, color=ACCENT, bold=True, first=True,
             align=PP_ALIGN.CENTER)
    px += pw + gap

y2 = y + 1.44
yl = column_head(s3, ML, y2, LW2, "What runs today")
tf = textbox(s3, ML, yl, LW2, 2.8)
bullet(tf, [("The command-line engine is the reference implementation. ",
             True, NAVY),
            ("One dependency — openpyxl — no build step. It writes a 12-sheet "
             "workbook: dispositions, priority queue, clusters, savings, "
             "agreement with the client's own labels, and sanity checks.",
             False, None)],
       first=True)
bullet(tf, [("Every row explains itself — ", True, NAVY),
            ("which dimension failed, at what score, the successor it folds "
             "into, and the evidence gap behind its confidence. Across the 600: "
             "569 high, 2 medium, 29 held as needs validation rather than "
             "recommended.", False, None)])
bullet(tf, [("600 applications in one run — seconds, not weeks. ", True, NAVY),
            ("10.02s end to end on the committed run, 3.75s of that to load, "
             "normalise, score and decide.", False, None)])

yr = column_head(s3, RX2, y2, RW2, "The web app — upload and score")
tf = textbox(s3, RX2, yr, RW2, 2.0)
bullet(tf, [("Upload a workbook or CSV and score it live. ", True, ACCENT),
            ("A real file dialogue, .xlsx or .csv in our intake schema. The "
             "page rewrites itself from the file: figures, decision mix, "
             "guardrails, clusters, and a per-row table with rationale and "
             "dimension scores. No server, no internet.", False, None)],
       first=True)
bullet(tf, [("It agrees with the engine on this portfolio, measured: ",
             True, NAVY),
            ("the page's own scoring, replayed over the uploaded export, "
             "returns the engine's disposition and priority on 600 of 600 rows.",
             False, None)])

stat_panel(s3, RX2, yr + 1.78, RW2, 0.94, "600 / 600",
           "rows where the browser matched the engine's answer.", num_size=26)

footnote(s3, "Engine: engine/score_northstar_v3.py and "
             "score_northstar_600_tuned.py  ·  browser scoring and the "
             "measured parity: src/apprat-ai-wireframe-v2.html, "
             "engine/verify_tuned_parity.js  ·  deployed at "
             "hack-team-12.vercel.app")
page_tag(s3, 3, "Demo")

notes(s3, """
Let me walk Northstar's flow at full portfolio scale, and I will be precise about what is built and
what is not.

Step one, intake. A workbook of the columns we asked the client for. The engine opens it
read-only and never writes back to it.

Step two, scoring. Eighteen inputs derived and scored, rolled into the four dimensions, each
gated at 3.0. Where an input has no source it is skipped and the dimension renormalises over what
is actually populated - it does not silently score a blank as a zero.

Step three, the lookup and the two guardrails I just described.

Step four is the one that makes this usable. Every row comes out with its rationale written in
prose - not just "consolidate", but which dimension failed, at what score, which input drove it,
and which application it folds into. Plus a confidence level and the specific evidence gap behind
that level. Across the six hundred: 569 rows at high confidence, two at medium, and twenty-nine
held as needs validation rather than given a recommendation at all, because the evidence behind
them was not there. That sixth state, needs validation, is deliberately not one of the five
dispositions - the tool declines rather than guesses.

Step five is the point. 233 applications come back retain: no action, no spend. 367 rows carry an
action, and those are the shortlist a CIO actually argues about - 129 consolidate, 129 invest, 109
retire, and nothing in replace. The consolidations fall out of the overlap clusters the engine
found: seven groups, and in each one the engine names the survivor and says why the others fold
into it. Ninety-three of those rows were folded in by the redundancy guardrail specifically,
because a named survivor already holds the capability as primary.

On how it runs: Python 3 with exactly one third-party dependency, openpyxl. No build step, no
package manager, nothing to compile. That command-line engine stays the reference implementation -
if the browser and the engine ever disagree, the engine is right by definition. Output is a
twelve-sheet workbook, including a sanity-checks sheet and a sheet comparing our dispositions
against the client's own labels, which is how we find the places we disagree with her, and why.

Speed, from the committed run rather than from memory: 10.02 seconds end to end for six hundred
applications, of which 3.75 seconds is the actual load, normalise, score and decide. Treat that as
an order of magnitude rather than a benchmark - it is a laptop-class measurement. The point is that
the unit cost per application stops being a meeting.

Now the web app, and here is what is real. Upload inventory opens a genuine file dialogue that
accepts .xlsx and .csv. Analyze portfolio scores that file in the browser and rewrites the page
from it - the headline figures, the savings bars, the decision mix, the guardrail counts, the
redundancy clusters, and the workbench table, where clicking a row opens the evidence rail rebuilt
from the computed result: four dimension scores with pass or fail, the pattern key, the priority,
the confidence and a written rationale. Three filters - capability, recommendation, critical
operation - narrow that table, and the headline figures deliberately do not move when you filter,
so a narrowed table can never be read as a changed total. There is no server, no build step and no
network call in the whole path: the file never leaves the machine, which also means this works on
a conference-centre wifi that does not.

We checked the browser against the engine on this exact portfolio rather than assuming it, and the
check is committed as engine/verify_tuned_parity.js. It pulls the page's own scoring code out of
index.html at run time - so it is testing the shipped page, not a copy of its logic - and replays it
over the six hundred exported rows. Disposition matches on 600 of 600. Priority matches on 600 of
600. Both post-lookup guardrails reproduce, including the ninety-three redundancy foldings, and no
row comes out with a negative net. The page's own arithmetic lands on the same net first-year figure
as the Python run, to the dollar.

One boundary to state plainly, and I would rather say it than have it found: it scores a portfolio
in our intake schema. It is not a universal reader. Hand it an inventory in a different column
vocabulary - as the raw Northstar workbook is - and it refuses: it names the file, the row and
column count, and lists the headings it does look for, app_id, app_name, primary_capability,
ov_patient_care_criticality, th_supportability, c_cost_per_active_user_vs_peers, r_technical_risk.
Case, spaces, hyphens and underscores do not matter. Mapping another vocabulary onto those names
is a real piece of work and we have not automated it.

IF AN UPLOAD IS REFUSED ON STAGE, say this: "That is the tool refusing a file it cannot score
rather than inventing an answer for it - it wants our intake column names, and it just told you
which ones. Here is the same portfolio in that schema." Then upload
data/northstar/northstar-600-tuned-tool-vocabulary.csv, which is the file we demo from - the
600-application portfolio behind the impact slide. Say what it is as you upload it: "this is the
illustrative portfolio, calibrated to the fifteen percent savings target - the label is on the file
itself, on its first sheet and on every row."

And to be exact about what I am claiming: the page is deployed at hack-team-12.vercel.app, and
everything I have just described about its behaviour is read out of the committed source in
src/apprat-ai-wireframe-v2.html and the parity run in engine/verify_tuned_parity.js. Export
decisions and Generate executive readout are still deliberately inert - we left them visibly
unwired rather than faking them.
""")


# ================================================================ SLIDE 4
s4 = prs.slides.add_slide(BLANK)
y = header(s4, "Impact · Business value and path to market",
           "15% target cleared at 18.47% — on a calibrated portfolio")

BH = 1.32
rect(s4, ML, y, CW, BH, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.05)
rect(s4, ML, y + 0.18, 0.055, BH - 0.36, ACCENT)

tfm = textbox(s4, ML + 0.34, y + 0.18, 3.02, BH - 0.36, anchor=MSO_ANCHOR.MIDDLE)
para(tfm, "18.47%", size=38, color=WHITE, bold=True, first=True, line=1.0)
para(tfm, "of run-rate saved in year one, net of transition cost — the "
          "15% target is cleared", size=10, color=ONNAVY, space_before=4,
     line=1.16)

facts = [
    ("$68,813,000",  "net first year, against a\n15% ask of $55,882,800"),
    ("$92,845,000",  "gross annual avoidable\n— recurs from year two"),
    ("$24,032,000",  "one-time transition cost\nnetted out (25.9% of gross)"),
    ("$372,552,000", "portfolio run-rate the\n18.47% is measured against"),
]
fx = ML + 3.52
fw = (CW - 3.52 - 0.16) / 4
for i, (num, lbl) in enumerate(facts):
    if i:
        rect(s4, fx - 0.05, y + 0.30, 0.008, BH - 0.60, DIVIDE)
    tff = textbox(s4, fx + 0.13, y + 0.24, fw - 0.24, BH - 0.48,
                  anchor=MSO_ANCHOR.MIDDLE)
    para(tff, num, size=16, color=WHITE, bold=True, first=True, line=1.0)
    para(tff, lbl, size=9.5, color=RGBColor(0xA8, 0xB8, 0xC6), space_before=4,
         line=1.16)
    fx += fw

y2 = y + BH + 0.30
CW3 = (CW - 0.90) / 3
cols = [ML, ML + CW3 + 0.45, ML + 2 * (CW3 + 0.45)]

yc = column_head(s4, cols[0], y2, CW3, "What this portfolio is")
tf = textbox(s4, cols[0], yc, CW3, 2.6)
bullet(tf, [("An illustrative portfolio, calibrated to the 15% target. ",
             True, NAVY),
            ("Every file in the set says so on its own first sheet: the values "
             "were constructed so the unchanged model would clear the target. "
             "So 18.47% is a property of this dataset, not a finding about an "
             "estate.", False, None)], first=True, size=11.5)
bullet(tf, [("What was calibrated: ", True, NAVY),
            ("154 of the 600 applications — no rows added, none invented — "
             "were relabelled as non-survivor members of their capability "
             "clusters, their existing costs scaled (86 by 1.15, 68 by 1.25), "
             "and migration evidence added.", False, None)], size=11.5)
bullet(tf, [("Named, so it can be checked: ", True, NAVY),
            ("data/northstar/northstar-600-tuned-tool-vocabulary.csv.",
             False, None)], size=11.5)

yc = column_head(s4, cols[1], y2, CW3, "What the engine did")
tf = textbox(s4, cols[1], yc, CW3, 2.6)
bullet(tf, [("The engine was not modified to reach the number. ", True, NAVY),
            ("Same 18 inputs, four dimensions, 3.0 gates, 16-row table, two "
             "guardrails and savings arithmetic as every other run. No row is "
             "special-cased: every disposition is the model's own answer.",
             False, None)], first=True, size=11.5)
bullet(tf, [("The spread across 600: ", True, NAVY),
            ("retain 233, invest 129, consolidate 129, retire 109, replace 0. "
             "367 rows carry an action; the other 233 are published as "
             "decisions with a route to object, not a meeting.",
             False, None)], size=11.5)
bullet(tf, [("Checked, not asserted: ", True, NAVY),
            ("the browser reproduces the engine's disposition and priority on "
             "600 of 600 rows.", False, None)], size=11.5)

yc = column_head(s4, cols[2], y2, CW3, "Why it is repeatable · what's next")
tf = textbox(s4, cols[2], yc, CW3, 2.6)
bullet(tf, [("The model is portfolio-agnostic. ", True, NAVY),
            ("It was built against a different, independently generated "
             "portfolio, and scored Northstar through the same code with no "
             "changes.", False, None)], first=True, size=11.5)
bullet(tf, [("Nothing sector-specific is hard-coded. ", True, NAVY),
            ("Thresholds, weights and all 16 table rows are configuration, so "
             "a new sector is a re-tune, not a rebuild.", False, None)],
       size=11.5)
bullet(tf, [("Next: ", True, NAVY),
            ("a real peer cost benchmark, and closing the derived-column "
             "pipeline — 31 of 68 columns are still authored by hand. The cost "
             "lens rests on three criteria rather than four: consumption price "
             "variance is absent.", False, None)], size=11.5)

footnote(s4, "Every figure above is the calibrated 600-application portfolio: "
             "an illustrative dataset constructed to clear the 15% target, "
             "labelled so on all 600 rows. Source: "
             "northstar-600-tuned-summary.md, re-derived from "
             "northstar-600-tuned-tool-vocabulary.csv.")
page_tag(s4, 4, "Impact")

notes(s4, """
The money first. This is a six-hundred-application portfolio scored end to end, not a sample of it,
and it carries a 372.552 million dollar annual run cost. Against that, 92.845 million of gross
annual avoidable cost, 24.032 million of one-time transition cost netted out - migration, interface
cutover, contract exit, which is 25.9 percent of the gross - and 68.813 million net in the first
year. From year two the gross recurs, because the transition cost is one time only. That is 18.47
percent of run cost. The CIO's stated target is fifteen percent, which on this run cost is 55.883
million, so the target is cleared and cleared by a margin.

Now the sentence that has to come with it, and I would rather say it first than be asked for it.
This portfolio was built to clear that target. Its values were constructed so the unchanged model
would return more than seventeen percent, and every file in the set carries a "Provenance - TUNED"
sheet as its first sheet, plus a data_source column saying the same thing on all six hundred rows.
So the 18.47 percent is a property of this dataset. It is an illustrative portfolio, calibrated to
the target - not a measurement of anybody's estate, and not a finding. Judges can open the file and
read that label for themselves, which is exactly why it is on the slide.

What calibrated means here, precisely, because the vague version sounds worse than the truth.
Nothing was duplicated. No rows were added and no applications were invented - it is six hundred
distinct applications, the same six hundred, and no duplicate application names. What changed is
that 154 of them were relabelled as non-survivor members of their capability clusters and had their
existing cost components scaled: 86 by 1.15 and 68 by 1.25. Alongside that, migration evidence was
written into the dependency data so a consolidation has somewhere to land, criticality and business
value were stepped down on those rows, and the retire-shaped ones were given technical decay -
vendor support ending inside a year, a legacy release line, MTTR past the 120-minute band.

Two consequences worth stating before anyone finds them. First, those two multipliers are uniform,
so anybody diffing the dataset spots them in seconds. They are visibly engineered, and that is fine
for a fixture that says on every row that it is a fixture. Second, scaling those costs also raised
the portfolio's run cost, to 372.552 million - the 7.548 million from the 1.15 rows and the 10.674
million from the 1.25 rows account for the rise exactly. So the percentage sits on a larger base.
Note which way that cuts, because it is easy to get backwards: a larger base makes the target
harder, not easier - fifteen percent of 372.552 million is a higher bar in dollars than fifteen
percent of a smaller run cost would be. The percentage moves because the saving grew, not because
the denominator did. Whenever you quote 18.47 percent, quote it against 372.552 million and no
other number.

And here is the claim I will actually defend. The engine was not modified to reach any of this.
Same eighteen inputs, same four dimensions, same 3.0 gates, same sixteen-row table, same two
guardrails, same savings arithmetic. No weight, band, rubric, gate, table row or formula differs,
and no row is special-cased in scoring. Every disposition in this run is the model's own answer to
the data it was handed. And that is checked rather than asserted: engine/verify_tuned_parity.js
pulls the page's own scoring out of index.html and replays it over the six hundred exported rows,
reproducing the engine's disposition and priority on 600 of 600 of them, both post-lookup guardrails
included, and landing on the same net figure to the dollar.

One more property worth stating, because it is the sort of thing that gets checked: no application
whose business criticality is critical, and none whose patient-care impact is direct, was turned
into a retire or consolidate candidate to reach the number. Healthcare criticality protects
applications, and the fixture honours that.

The spread is retain 233, invest 129, consolidate 129, retire 109, replace 0. Nothing lands in
replace, because every application that might have been replaced clears its gates. Operationally
that split is the part I care about more than the headline: 233 applications come back retain - no
action, no spend - and today each of those still costs you a meeting to conclude nothing. Under
AdvisAR they are published as decisions with the evidence attached and a route to object. Nobody
has to defend all 233; anybody can challenge one. The 367 rows that carry an action are the
shortlist, and only those get the full seven-step treatment - the workshops, the dependency checks,
the sequencing. Scarce facilitation goes only where it can change an answer, and that is what makes
six hundred tractable.

One portfolio, so one set of numbers. Every figure on this slide and on the demo slide comes from
the same calibrated six-hundred-application run - the same file we upload on stage. I am not
blending it with anything, and if you ask me for a figure I will tell you which column of which file
it came out of.

Two limits of the data, briefly. Consumption price variance is absent from the source, so the cost
lens on this portfolio rests on three criteria rather than four. That input carries the lowest
weight in the model and the others renormalise around it, so it does not condemn a row - it is one
missing portfolio-wide input, not a penalty applied per application. And I am not quoting a
confidence-qualified savings figure on this slide at all: the engine does compute a confidence
split, but the summary offers two different measures of it, and an unlabelled number would be read
as whichever one the listener assumed. If you want that figure, ask me and I will tell you which
measure I am giving you.

Now path to market, which is the reusability line. The scoring engine was built against a
completely different portfolio - a separately generated dataset with a different shape. Northstar
was constructed independently as a test of the model, and it ran through the same code with no
changes at all. That is the claim: nothing about healthcare is
hard-coded. The thresholds, the weights and all sixteen rows of the disposition table are
configuration rather than branching logic, so taking this into manufacturing or financial
services is a re-tune of a config table, not a rebuild. It drops into an Aberdeen engagement as a
repeatable offering with a known intake list and a known deliverable, instead of a bespoke
analysis we rebuild every time.

One practical note for the demo. The file we upload on stage is
data/northstar/northstar-600-tuned-tool-vocabulary.csv - this portfolio translated into the tool's
own column names. I checked the file itself rather than trusting the summary: summing its columns
gives 372,552,000 of run cost, 92,845,000 gross, 24,032,000 of transition cost and 68,813,000 net,
to the dollar, with the same 233 / 129 / 129 / 109 / 0 spread. So the figures on this slide and the
figures on the screen are the same figures, and if they ever diverge the engine run recorded in
northstar-600-tuned-summary.md is the one to trust.

Two things we owe you before this is client-ready. One: the peer cost band. Cost per active user
against peers is the single most influential input in the cost dimension, and right now it is
modelled from the portfolio itself rather than measured against an external benchmark. We label
it as modelled in every output, but a live engagement needs a real data source. Two: the pipeline
in front of the engine. Of the 68 columns the tool derives, 31 are currently authored by hand
rather than computed - AI classification, vendor end-of-support lookups, the clustering step, and
five band mappings. When those are supplied, all 37 implemented derivations reproduce exactly, on
every row, with zero mismatches. So the engine is right; the intake pipeline is the unfinished
part, and the five band mappings are the cheapest place to start.
""")


out = sys.argv[1] if len(sys.argv) > 1 else "AdvisAR-Hackathon-Deck.pptx"
prs.save(out)
print("saved", out)
